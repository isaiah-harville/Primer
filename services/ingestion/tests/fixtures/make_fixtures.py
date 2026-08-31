"""Generate the PDF fixtures.

Checked in as a script so the binaries beside it are reviewable: a fixture
nobody can regenerate is a fixture nobody can trust or change.

Run: uv run python services/ingestion/tests/fixtures/make_fixtures.py
"""

from __future__ import annotations

import zlib
from pathlib import Path

HERE = Path(__file__).parent


def build_pdf(objects: list[bytes]) -> bytes:
    """Assemble numbered objects into a PDF with a correct xref table."""
    out = bytearray(b"%PDF-1.7\n%\xe2\xe3\xcf\xd3\n")
    offsets = [0]
    for number, body in enumerate(objects, start=1):
        offsets.append(len(out))
        out += f"{number} 0 obj\n".encode() + body + b"\nendobj\n"

    start_xref = len(out)
    out += f"xref\n0 {len(objects) + 1}\n".encode()
    out += b"0000000000 65535 f \n"
    for offset in offsets[1:]:
        out += f"{offset:010d} 00000 n \n".encode()
    out += (
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{start_xref}\n%%EOF\n"
    ).encode()
    return bytes(out)


def stream_object(data: bytes) -> bytes:
    return b"<< /Length " + str(len(data)).encode() + b" >>\nstream\n" + data + b"\nendstream"


def text_paper() -> bytes:
    """One page of real, extractable text: a heading and two paragraphs.

    Single page on purpose. Docling's PDF backend rejects the second page of
    anything this generator emits, whatever the object layout, and chasing
    that would be reverse-engineering a parser rather than testing Primer.
    One page still proves page numbers are extracted and attached; it does
    not prove they increment, which is the coverage this trades away.
    """
    content = (
        b"BT /F1 18 Tf 72 720 Td (Retrieval Augmented Generation) Tj ET\n"
        b"BT /F1 11 Tf 72 690 Td (Grounding answers in cited sources reduces "
        b"unsupported claims.) Tj ET\n"
        b"BT /F1 11 Tf 72 660 Td (Recall at rank ten was the decisive metric "
        b"for this corpus.) Tj ET\n"
    )
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R "
        b"/Resources << /Font << /F1 5 0 R >> >> >>",
        stream_object(content),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    return build_pdf(objects)


def scanned_paper() -> bytes:
    """A page that is only an image: what a scanner produces, and what OCR is for."""
    width = height = 8
    raw = bytes([0x80]) * (width * height)
    image = zlib.compress(raw)
    content = b"q 612 0 0 792 0 0 cm /Im0 Do Q\n"
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R "
        b"/Resources << /XObject << /Im0 5 0 R >> >> >>",
        stream_object(content),
        b"<< /Type /XObject /Subtype /Image /Width "
        + str(width).encode()
        + b" /Height "
        + str(height).encode()
        + b" /ColorSpace /DeviceGray /BitsPerComponent 8 /Filter /FlateDecode /Length "
        + str(len(image)).encode()
        + b" >>\nstream\n"
        + image
        + b"\nendstream",
    ]
    return build_pdf(objects)


def scanned_with_text() -> bytes:
    """A page that is a picture of words: only OCR can read it.

    The text is rendered to an image and embedded as a JPEG, so the PDF
    contains no character codes at all. A parser without OCR finds nothing
    here, which is precisely what makes it a test of OCR rather than of PDF
    text extraction.
    """
    from io import BytesIO

    from PIL import Image, ImageDraw, ImageFont

    image = Image.new("RGB", (1240, 400), "white")
    draw = ImageDraw.Draw(image)
    draw.text((60, 140), "SCANNED EVIDENCE", fill="black", font=ImageFont.load_default(size=90))

    buffer = BytesIO()
    image.save(buffer, format="JPEG", quality=95)
    jpeg = buffer.getvalue()

    content = b"q 612 0 0 198 0 400 cm /Im0 Do Q\n"
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R "
        b"/Resources << /XObject << /Im0 5 0 R >> >> >>",
        stream_object(content),
        b"<< /Type /XObject /Subtype /Image /Width "
        + str(image.width).encode()
        + b" /Height "
        + str(image.height).encode()
        + b" /ColorSpace /DeviceRGB /BitsPerComponent 8 /Filter /DCTDecode /Length "
        + str(len(jpeg)).encode()
        + b" >>\nstream\n"
        + jpeg
        + b"\nendstream",
    ]
    return build_pdf(objects)


def slide_deck() -> bytes:
    """A two-slide deck with titles and body text.

    Built with python-pptx rather than by hand: a .pptx is a ZIP of related
    OOXML parts, and hand-assembling one would be testing the fixture
    generator rather than the parser.
    """
    from io import BytesIO

    from pptx import Presentation

    presentation = Presentation()
    title_and_body = presentation.slide_layouts[1]

    first = presentation.slides.add_slide(title_and_body)
    first.shapes.title.text = "Retrieval Augmented Generation"
    first.placeholders[1].text = "Grounding answers in cited sources reduces unsupported claims."

    second = presentation.slides.add_slide(title_and_body)
    second.shapes.title.text = "Evaluation"
    second.placeholders[1].text = "Recall at rank ten was the decisive metric for this corpus."

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


if __name__ == "__main__":
    (HERE / "text-paper.pdf").write_bytes(text_paper())
    (HERE / "scanned-paper.pdf").write_bytes(scanned_paper())
    (HERE / "scanned-with-text.pdf").write_bytes(scanned_with_text())
    (HERE / "slides.pptx").write_bytes(slide_deck())
    print("wrote text-paper.pdf, scanned-paper.pdf, scanned-with-text.pdf, and slides.pptx")
